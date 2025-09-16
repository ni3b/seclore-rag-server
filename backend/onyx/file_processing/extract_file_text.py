import io
import json
import os
import re
import zipfile
import base64
from collections.abc import Callable
from collections.abc import Iterator
from email.parser import Parser as EmailParser
from io import BytesIO
from pathlib import Path
from typing import Any
from typing import Dict
from typing import IO

import chardet
import docx  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from docx import Document
from fastapi import UploadFile
from pypdf import PdfReader
from pypdf.errors import PdfStreamError

from onyx.configs.constants import DANSWER_METADATA_FILENAME
from onyx.configs.constants import FileOrigin
from onyx.file_processing.html_utils import parse_html_page_basic
from onyx.file_processing.unstructured import get_unstructured_api_key
from onyx.file_processing.unstructured import unstructured_to_text
from onyx.file_store.file_store import FileStore
from onyx.utils.logger import setup_logger

logger = setup_logger()


TEXT_SECTION_SEPARATOR = "\n\n"


PLAIN_TEXT_FILE_EXTENSIONS = [
    ".txt",
    ".md",
    ".mdx",
    ".conf",
    ".log",
    ".json",
    ".csv",
    ".tsv",
    ".xml",
    ".yml",
    ".yaml",
]

# Image file extensions for OCR processing
IMAGE_FILE_EXTENSIONS = [
    ".jpg",
    ".jpeg",
    ".png",
    ".gif",
    ".bmp",
    ".tiff",
    ".tif",
    ".webp",
    ".svg",
    ".avif",
    ".heic",
    ".heif",
]

VALID_FILE_EXTENSIONS = PLAIN_TEXT_FILE_EXTENSIONS + [
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".eml",
    ".epub",
    ".html",
] + IMAGE_FILE_EXTENSIONS


def is_text_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in PLAIN_TEXT_FILE_EXTENSIONS)


def is_image_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in IMAGE_FILE_EXTENSIONS)


def get_file_ext(file_path_or_name: str | Path) -> str:
    _, extension = os.path.splitext(file_path_or_name)
    # standardize all extensions to be lowercase so that checks against
    # VALID_FILE_EXTENSIONS and similar will work as intended
    return extension.lower()


def is_valid_file_ext(ext: str) -> bool:
    return ext in VALID_FILE_EXTENSIONS


def is_text_file(file: IO[bytes]) -> bool:
    """
    checks if the first 1024 bytes only contain printable or whitespace characters
    if it does, then we say its a plaintext file
    """
    raw_data = file.read(1024)
    text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7F})
    return all(c in text_chars for c in raw_data)


def detect_encoding(file: IO[bytes]) -> str:
    raw_data = file.read(50000)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    file.seek(0)
    return encoding


def is_macos_resource_fork_file(file_name: str) -> bool:
    return os.path.basename(file_name).startswith("._") and file_name.startswith(
        "__MACOSX"
    )


# To include additional metadata in the search index, add a .onyx_metadata.json file
# to the zip file. This file should contain a list of objects with the following format:
# [{ "filename": "file1.txt", "link": "https://example.com/file1.txt" }]
def load_files_from_zip(
    zip_file_io: IO,
    ignore_macos_resource_fork_files: bool = True,
    ignore_dirs: bool = True,
) -> Iterator[tuple[zipfile.ZipInfo, IO[Any], dict[str, Any]]]:
    with zipfile.ZipFile(zip_file_io, "r") as zip_file:
        zip_metadata = {}
        try:
            metadata_file_info = zip_file.getinfo(DANSWER_METADATA_FILENAME)
            with zip_file.open(metadata_file_info, "r") as metadata_file:
                try:
                    zip_metadata = json.load(metadata_file)
                    if isinstance(zip_metadata, list):
                        # convert list of dicts to dict of dicts
                        zip_metadata = {d["filename"]: d for d in zip_metadata}
                except json.JSONDecodeError:
                    logger.warn(f"Unable to load {DANSWER_METADATA_FILENAME}")
        except KeyError:
            logger.info(f"No {DANSWER_METADATA_FILENAME} file")

        for file_info in zip_file.infolist():
            with zip_file.open(file_info.filename, "r") as file:
                if ignore_dirs and file_info.is_dir():
                    continue

                if (
                    ignore_macos_resource_fork_files
                    and is_macos_resource_fork_file(file_info.filename)
                ) or file_info.filename == DANSWER_METADATA_FILENAME:
                    continue
                yield file_info, file, zip_metadata.get(file_info.filename, {})


def _extract_onyx_metadata(line: str) -> dict | None:
    html_comment_pattern = r"<!--\s*DANSWER_METADATA=\{(.*?)\}\s*-->"
    hashtag_pattern = r"#DANSWER_METADATA=\{(.*?)\}"

    html_comment_match = re.search(html_comment_pattern, line)
    hashtag_match = re.search(hashtag_pattern, line)

    if html_comment_match:
        json_str = html_comment_match.group(1)
    elif hashtag_match:
        json_str = hashtag_match.group(1)
    else:
        return None

    try:
        return json.loads("{" + json_str + "}")
    except json.JSONDecodeError:
        return None


def read_text_file(
    file: IO,
    encoding: str = "utf-8",
    errors: str = "replace",
    ignore_onyx_metadata: bool = True,
) -> tuple[str, dict]:
    metadata = {}
    file_content_raw = ""
    for ind, line in enumerate(file):
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = (
                line.decode(encoding, errors=errors)
                if isinstance(line, bytes)
                else line
            )

        if ind == 0:
            metadata_or_none = (
                None if ignore_onyx_metadata else _extract_onyx_metadata(line)
            )
            if metadata_or_none is not None:
                metadata = metadata_or_none
            else:
                file_content_raw += line
        else:
            file_content_raw += line

    return file_content_raw, metadata


def pdf_to_text(file: IO[Any], pdf_pass: str | None = None) -> str:
    """Extract text from a PDF file."""
    # Return only the extracted text from read_pdf_file
    text, _ = read_pdf_file(file, pdf_pass)
    return text


def read_pdf_file(
    file: IO[Any],
    pdf_pass: str | None = None,
) -> tuple[str, dict]:
    metadata: Dict[str, Any] = {}
    try:
        pdf_reader = PdfReader(file)

        # If marked as encrypted and a password is provided, try to decrypt
        if pdf_reader.is_encrypted and pdf_pass is not None:
            decrypt_success = False
            if pdf_pass is not None:
                try:
                    decrypt_success = pdf_reader.decrypt(pdf_pass) != 0
                except Exception:
                    logger.error("Unable to decrypt pdf")

            if not decrypt_success:
                # By user request, keep files that are unreadable just so they
                # can be discoverable by title.
                return "", metadata
        elif pdf_reader.is_encrypted:
            logger.warning("No Password available to decrypt pdf, returning empty")
            return "", metadata

        # Extract metadata from the PDF, removing leading '/' from keys if present
        # This standardizes the metadata keys for consistency
        metadata = {}
        if pdf_reader.metadata is not None:
            for key, value in pdf_reader.metadata.items():
                clean_key = key.lstrip("/")
                if isinstance(value, str) and value.strip():
                    metadata[clean_key] = value

                elif isinstance(value, list) and all(
                    isinstance(item, str) for item in value
                ):
                    metadata[clean_key] = ", ".join(value)

        return (
            TEXT_SECTION_SEPARATOR.join(
                page.extract_text() for page in pdf_reader.pages
            ),
            metadata,
        )
    except PdfStreamError:
        logger.exception("PDF file is not a valid PDF")
    except Exception:
        logger.exception("Failed to read PDF")

    # File is still discoverable by title
    # but the contents are not included as they cannot be parsed
    return "", metadata


def docx_to_text(file: IO[Any]) -> str:
    def is_simple_table(table: docx.table.Table) -> bool:
        for row in table.rows:
            # No omitted cells
            if row.grid_cols_before > 0 or row.grid_cols_after > 0:
                return False

            # No nested tables
            if any(cell.tables for cell in row.cells):
                return False

        return True

    def extract_cell_text(cell: docx.table._Cell) -> str:
        cell_paragraphs = [para.text.strip() for para in cell.paragraphs]
        return " ".join(p for p in cell_paragraphs if p) or "N/A"

    paragraphs = []
    doc = docx.Document(file)
    for item in doc.iter_inner_content():
        if isinstance(item, docx.text.paragraph.Paragraph):
            paragraphs.append(item.text)

        elif isinstance(item, docx.table.Table):
            if not item.rows or not is_simple_table(item):
                continue

            # Every row is a new line, joined with a single newline
            table_content = "\n".join(
                [
                    ",\t".join(extract_cell_text(cell) for cell in row.cells)
                    for row in item.rows
                ]
            )
            paragraphs.append(table_content)

    # Docx already has good spacing between paragraphs
    return "\n".join(paragraphs)


def pptx_to_text(file: IO[Any]) -> str:
    presentation = pptx.Presentation(file)
    text_content = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        extracted_text = f"\nSlide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
        text_content.append(extracted_text)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def xlsx_to_text(file: IO[Any]) -> str:
    workbook = openpyxl.load_workbook(file, read_only=True)
    text_content = []
    for sheet in workbook.worksheets:
        sheet_string = "\n".join(
            ",".join(map(str, row))
            for row in sheet.iter_rows(min_row=1, values_only=True)
        )
        text_content.append(sheet_string)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def eml_to_text(file: IO[Any]) -> str:
    text_file = io.TextIOWrapper(file, encoding=detect_encoding(file))
    parser = EmailParser()
    message = parser.parse(text_file)
    text_content = []
    for part in message.walk():
        if part.get_content_type().startswith("text/plain"):
            text_content.append(part.get_payload())
    return TEXT_SECTION_SEPARATOR.join(text_content)


def epub_to_text(file: IO[Any]) -> str:
    with zipfile.ZipFile(file) as epub:
        text_content = []
        for item in epub.infolist():
            if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                with epub.open(item) as html_file:
                    text_content.append(parse_html_page_basic(html_file))
        return TEXT_SECTION_SEPARATOR.join(text_content)


def file_io_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    file_content_raw, _ = read_text_file(file, encoding=encoding)
    return file_content_raw


def image_to_text(file: IO[Any]) -> str:
    """Extract text from image files using OCR via Unstructured API."""
    try:
        # Reset file pointer to beginning
        file.seek(0)
        
        # Use the existing unstructured_to_text function which supports OCR for images
        if get_unstructured_api_key():
            # The unstructured API supports OCR for images, including:
            # - Text extraction from scanned documents
            # - OCR from images (JPG, PNG, GIF, BMP, TIFF, WebP, etc.)
            # - Vision-based text recognition
            from onyx.file_processing.unstructured import unstructured_to_text
            return unstructured_to_text(file, "image_file")
        else:
            # Fallback: attempt basic text extraction if no API key
            logger.warning("No Unstructured API key found. OCR functionality requires Unstructured API access.")
            return ""
            
    except Exception as e:
        logger.warning(f"Failed to extract text from image: {str(e)}")
        return ""


def _generate_image_description_with_sonnet(file: IO[Any], file_name: str) -> str:
    """
    Generate a description of the image using Claude Sonnet 4 vision capabilities.
    """
    try:
        # Read image data and encode as base64
        file.seek(0)
        image_data = file.read()
        base64_image = base64.b64encode(image_data).decode('utf-8')
        
        # Get the default LLM which should be Claude Sonnet 4 with vision
        from onyx.llm.factory import get_default_llm
        
        llm = get_default_llm()
        
        # Create a vision prompt for image description
        vision_prompt = """Please provide a detailed description of this image. Include:
1. What objects, people, or scenes are visible
2. The setting or context
3. Any text visible in the image (if any)
4. Colors, composition, and notable visual elements
5. Any relevant details that would help someone understand the content

Provide a comprehensive but concise description that would be useful for search and retrieval."""

        try:
            # Use Claude Sonnet's vision capabilities
            # Create message with image
            from onyx.llm.utils import build_content_with_imgs
            
            # Create content with image
            content = build_content_with_imgs(
                message=vision_prompt,
                b64_imgs=[base64_image]
            )
            
            # Invoke the LLM with the image content
            response = llm.invoke(content)
            return response.strip() if response else ""
            
        except Exception as e:
            logger.warning(f"Claude Sonnet vision call failed: {str(e)}")
            
            # Fallback: try with Unstructured API if it supports image descriptions
            if get_unstructured_api_key():
                try:
                    file.seek(0)
                    from onyx.file_processing.unstructured import unstructured_to_text
                    # Some versions of Unstructured API can provide image descriptions
                    description = unstructured_to_text(file, file_name)
                    if description and len(description) > 10:  # Basic check for meaningful description
                        return description
                except:
                    pass
            
            return ""
            
    except Exception as e:
        logger.warning(f"Image description generation failed: {str(e)}")
        return ""


def process_image_comprehensive(file: IO[Any], file_name: str = "image") -> tuple[str, dict[str, Any]]:
    """
    Comprehensive image processing that extracts:
    1. OCR text from the image
    2. Visual description/summary of the image using Claude Sonnet 4
    3. Metadata for image indexing
    
    Returns:
        tuple: (combined_text, metadata_dict)
    """
    try:
        # Reset file pointer to beginning
        file.seek(0)
        
        # Extract OCR text
        ocr_text = ""
        image_description = ""
        
        if get_unstructured_api_key():
            try:
                from onyx.file_processing.unstructured import unstructured_to_text
                ocr_text = unstructured_to_text(file, file_name)
            except Exception as e:
                logger.warning(f"OCR extraction failed: {str(e)}")
        
        # Generate image description using Claude Sonnet 4 vision
        file.seek(0)
        try:
            image_description = _generate_image_description_with_sonnet(file, file_name)
        except Exception as e:
            logger.warning(f"Image description generation failed: {str(e)}")
        
        # Combine OCR text and image description
        text_parts = []
        if ocr_text.strip():
            text_parts.append(f"OCR Text: {ocr_text.strip()}")
        if image_description.strip():
            text_parts.append(f"Image Description: {image_description.strip()}")
        
        combined_text = "\n\n".join(text_parts) if text_parts else ""
        
        # Create comprehensive metadata
        metadata = {
            "file_type": "image",
            "has_ocr_text": bool(ocr_text.strip()),
            "has_description": bool(image_description.strip()),
            "ocr_text": ocr_text.strip(),
            "image_description": image_description.strip(),
            "processing_method": "comprehensive_image_processing",
            "vision_model": "claude-sonnet-4"
        }
        
        return combined_text, metadata
        
    except Exception as e:
        logger.error(f"Failed to process image comprehensively: {str(e)}")
        return "", {"file_type": "image", "error": str(e)}


def extract_file_text(
    file: IO[Any],
    file_name: str,
    break_on_unprocessable: bool = True,
    extension: str | None = None,
) -> str:
    extension_to_function: dict[str, Callable[[IO[Any]], str]] = {
        ".pdf": pdf_to_text,
        ".docx": docx_to_text,
        ".pptx": pptx_to_text,
        ".xlsx": xlsx_to_text,
        ".eml": eml_to_text,
        ".epub": epub_to_text,
        ".html": parse_html_page_basic,
        # Image file extensions for OCR processing
        ".jpg": image_to_text,
        ".jpeg": image_to_text,
        ".png": image_to_text,
        ".gif": image_to_text,
        ".bmp": image_to_text,
        ".tiff": image_to_text,
        ".tif": image_to_text,
        ".webp": image_to_text,
        ".svg": image_to_text,
        ".avif": image_to_text,
        ".heic": image_to_text,
        ".heif": image_to_text,
    }

    try:
        if get_unstructured_api_key():
            try:
                return unstructured_to_text(file, file_name)
            except Exception as unstructured_error:
                logger.error(
                    f"Failed to process with Unstructured: {str(unstructured_error)}. Falling back to normal processing."
                )
                # Fall through to normal processing

        if file_name or extension:
            if extension is not None:
                final_extension = extension
            elif file_name is not None:
                final_extension = get_file_ext(file_name)

            if is_valid_file_ext(final_extension):
                return extension_to_function.get(final_extension, file_io_to_text)(file)

        # Either the file somehow has no name or the extension is not one that we recognize
        if is_text_file(file):
            return file_io_to_text(file)

        raise ValueError("Unknown file extension and unknown text encoding")

    except Exception as e:
        if break_on_unprocessable:
            raise RuntimeError(
                f"Failed to process file {file_name or 'Unknown'}: {str(e)}"
            ) from e
        logger.warning(f"Failed to process file {file_name or 'Unknown'}: {str(e)}")
        return ""


def convert_docx_to_txt(
    file: UploadFile, file_store: FileStore, file_path: str
) -> None:
    file.file.seek(0)
    docx_content = file.file.read()
    doc = Document(BytesIO(docx_content))

    # Extract text from the document
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)

    # Join the extracted text
    text_content = "\n".join(full_text)

    txt_file_path = docx_to_txt_filename(file_path)
    file_store.save_file(
        file_name=txt_file_path,
        content=BytesIO(text_content.encode("utf-8")),
        display_name=file.filename,
        file_origin=FileOrigin.CONNECTOR,
        file_type="text/plain",
    )


def docx_to_txt_filename(file_path: str) -> str:
    """
    Convert a .docx file path to its corresponding .txt file path.
    """
    return file_path.rsplit(".", 1)[0] + ".txt"

import io
import json
import os
import re
import zipfile
from collections.abc import Callable
from collections.abc import Iterator
from collections.abc import Sequence
from email.parser import Parser as EmailParser
from enum import auto
from enum import IntFlag
from io import BytesIO
from pathlib import Path
from typing import Any
from typing import IO
from typing import NamedTuple
from zipfile import BadZipFile

import chardet
from markitdown import FileConversionException
from markitdown import MarkItDown
from markitdown import UnsupportedFormatException
from PIL import Image
from pypdf import PdfReader
from pypdf.errors import PdfStreamError

from onyx.configs.constants import ONYX_METADATA_FILENAME
from onyx.configs.llm_configs import get_image_extraction_and_analysis_enabled
from onyx.file_processing.file_validation import TEXT_MIME_TYPE
from onyx.file_processing.html_utils import parse_html_page_basic
from onyx.file_processing.unstructured import get_unstructured_api_key
from onyx.file_processing.unstructured import unstructured_to_text
from onyx.utils.logger import setup_logger

logger = setup_logger()

# NOTE(rkuo): Unify this with upload_files_for_chat and file_valiation.py
TEXT_SECTION_SEPARATOR = "\n\n"

ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS = [
    ".txt",
    ".md",
    ".mdx",
    ".conf",
    ".log",
    ".json",
    ".csv",
    ".tsv",
    ".xml",
    ".yml",
    ".yaml",
]

ACCEPTED_DOCUMENT_FILE_EXTENSIONS = [
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".eml",
    ".epub",
    ".html",
]

ACCEPTED_IMAGE_FILE_EXTENSIONS = [
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
]

ALL_ACCEPTED_FILE_EXTENSIONS = (
    ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS
    + ACCEPTED_DOCUMENT_FILE_EXTENSIONS
    + ACCEPTED_IMAGE_FILE_EXTENSIONS
)

IMAGE_MEDIA_TYPES = [
    "image/png",
    "image/jpeg",
    "image/webp",
]


class OnyxExtensionType(IntFlag):
    Plain = auto()
    Document = auto()
    Multimedia = auto()
    All = Plain | Document | Multimedia


def is_text_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS)


def get_file_ext(file_path_or_name: str | Path) -> str:
    _, extension = os.path.splitext(file_path_or_name)
    return extension.lower()


def is_valid_media_type(media_type: str) -> bool:
    return media_type in IMAGE_MEDIA_TYPES


def is_accepted_file_ext(ext: str, ext_type: OnyxExtensionType) -> bool:
    if ext_type & OnyxExtensionType.Plain:
        if ext in ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS:
            return True

    if ext_type & OnyxExtensionType.Document:
        if ext in ACCEPTED_DOCUMENT_FILE_EXTENSIONS:
            return True

    if ext_type & OnyxExtensionType.Multimedia:
        if ext in ACCEPTED_IMAGE_FILE_EXTENSIONS:
            return True

    return False


def is_text_file(file: IO[bytes]) -> bool:
    """
    checks if the first 1024 bytes only contain printable or whitespace characters
    if it does, then we say it's a plaintext file
    """
    raw_data = file.read(1024)
    file.seek(0)
    text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7F})
    return all(c in text_chars for c in raw_data)


def detect_encoding(file: IO[bytes]) -> str:
    raw_data = file.read(50000)
    file.seek(0)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    return encoding


def is_macos_resource_fork_file(file_name: str) -> bool:
    return os.path.basename(file_name).startswith("._") and file_name.startswith(
        "__MACOSX"
    )


def to_bytesio(stream: IO[bytes]) -> BytesIO:
    if isinstance(stream, BytesIO):
        return stream
    data = stream.read()  # consumes the stream!
    return BytesIO(data)


def load_files_from_zip(
    zip_file_io: IO,
    ignore_macos_resource_fork_files: bool = True,
    ignore_dirs: bool = True,
) -> Iterator[tuple[zipfile.ZipInfo, IO[Any]]]:
    """
    Iterates through files in a zip archive, yielding (ZipInfo, file handle) pairs.
    """
    with zipfile.ZipFile(zip_file_io, "r") as zip_file:
        for file_info in zip_file.infolist():
            if ignore_dirs and file_info.is_dir():
                continue

            if (
                ignore_macos_resource_fork_files
                and is_macos_resource_fork_file(file_info.filename)
            ) or file_info.filename == ONYX_METADATA_FILENAME:
                continue

            with zip_file.open(file_info.filename, "r") as subfile:
                # Try to match by exact filename first
                yield file_info, subfile


def _extract_onyx_metadata(line: str) -> dict | None:
    """
    Example: first line has:
        <!-- ONYX_METADATA={"title": "..."} -->
      or
        #ONYX_METADATA={"title":"..."}
    """
    html_comment_pattern = r"<!--\s*ONYX_METADATA=\{(.*?)\}\s*-->"
    hashtag_pattern = r"#ONYX_METADATA=\{(.*?)\}"

    html_comment_match = re.search(html_comment_pattern, line)
    hashtag_match = re.search(hashtag_pattern, line)

    if html_comment_match:
        json_str = html_comment_match.group(1)
    elif hashtag_match:
        json_str = hashtag_match.group(1)
    else:
        return None

    try:
        return json.loads("{" + json_str + "}")
    except json.JSONDecodeError:
        return None


def read_text_file(
    file: IO,
    encoding: str = "utf-8",
    errors: str = "replace",
    ignore_onyx_metadata: bool = True,
) -> tuple[str, dict]:
    """
    For plain text files. Optionally extracts Onyx metadata from the first line.
    """
    metadata = {}
    file_content_raw = ""
    for ind, line in enumerate(file):
        # decode
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = (
                line.decode(encoding, errors=errors)
                if isinstance(line, bytes)
                else line
            )

        # optionally parse metadata in the first line
        if ind == 0 and not ignore_onyx_metadata:
            potential_meta = _extract_onyx_metadata(line)
            if potential_meta is not None:
                metadata = potential_meta
                continue

        file_content_raw += line

    return file_content_raw, metadata


def pdf_to_text(file: IO[Any], pdf_pass: str | None = None) -> str:
    """
    Extract text from a PDF. For embedded images, a more complex approach is needed.
    This is a minimal approach returning text only.
    """
    text, _, _ = read_pdf_file(file, pdf_pass)
    return text


def read_pdf_file(
    file: IO[Any], pdf_pass: str | None = None, extract_images: bool = False
) -> tuple[str, dict[str, Any], Sequence[tuple[bytes, str]]]:
    """
    Returns the text, basic PDF metadata, and optionally extracted images.
    """
    metadata: dict[str, Any] = {}
    extracted_images: list[tuple[bytes, str]] = []
    try:
        pdf_reader = PdfReader(file)

        if pdf_reader.is_encrypted and pdf_pass is not None:
            decrypt_success = False
            try:
                decrypt_success = pdf_reader.decrypt(pdf_pass) != 0
            except Exception:
                logger.error("Unable to decrypt pdf")

            if not decrypt_success:
                return "", metadata, []
        elif pdf_reader.is_encrypted:
            logger.warning("No Password for an encrypted PDF, returning empty text.")
            return "", metadata, []

        # Basic PDF metadata
        if pdf_reader.metadata is not None:
            for key, value in pdf_reader.metadata.items():
                clean_key = key.lstrip("/")
                if isinstance(value, str) and value.strip():
                    metadata[clean_key] = value
                elif isinstance(value, list) and all(
                    isinstance(item, str) for item in value
                ):
                    metadata[clean_key] = ", ".join(value)

        text = TEXT_SECTION_SEPARATOR.join(
            page.extract_text() for page in pdf_reader.pages
        )

        if extract_images:
            for page_num, page in enumerate(pdf_reader.pages):
                for image_file_object in page.images:
                    image = Image.open(io.BytesIO(image_file_object.data))
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format=image.format)
                    img_bytes = img_byte_arr.getvalue()

                    image_name = (
                        f"page_{page_num + 1}_image_{image_file_object.name}."
                        f"{image.format.lower() if image.format else 'png'}"
                    )
                    extracted_images.append((img_bytes, image_name))

        return text, metadata, extracted_images

    except PdfStreamError:
        logger.exception("Invalid PDF file")
    except Exception:
        logger.exception("Failed to read PDF")

    return "", metadata, []


def extract_docx_images(docx_bytes: IO[Any]) -> list[tuple[bytes, str]]:
    """
    Given the bytes of a docx file, extract all the images.
    Returns a list of tuples (image_bytes, image_name).
    """
    out = []
    try:
        with zipfile.ZipFile(docx_bytes) as z:
            for name in z.namelist():
                if name.startswith("word/media/"):
                    out.append((z.read(name), name.split("/")[-1]))
    except Exception:
        logger.exception("Failed to extract all docx images")
    return out


def docx_to_text_and_images(
    file: IO[Any], file_name: str = ""
) -> tuple[str, Sequence[tuple[bytes, str]]]:
    """
    Extract text from a docx.
    Return (text_content, list_of_images).
    """
    md = MarkItDown(enable_plugins=False)
    try:
        doc = md.convert(to_bytesio(file))
    except (
        BadZipFile,
        ValueError,
        FileConversionException,
        UnsupportedFormatException,
    ) as e:
        logger.warning(
            f"Failed to extract docx {file_name or 'docx file'}: {e}. Attempting to read as text file."
        )

        # May be an invalid docx, but still a valid text file
        file.seek(0)
        encoding = detect_encoding(file)
        text_content_raw, _ = read_text_file(
            file, encoding=encoding, ignore_onyx_metadata=False
        )
        return text_content_raw or "", []

    file.seek(0)
    return doc.markdown, extract_docx_images(to_bytesio(file))


def pptx_to_text(file: IO[Any], file_name: str = "") -> str:
    md = MarkItDown(enable_plugins=False)
    try:
        presentation = md.convert(to_bytesio(file))
    except (
        BadZipFile,
        ValueError,
        FileConversionException,
        UnsupportedFormatException,
    ) as e:
        error_str = f"Failed to extract text from {file_name or 'pptx file'}: {e}"
        logger.warning(error_str)
        return ""
    return presentation.markdown


def xlsx_to_text(file: IO[Any], file_name: str = "") -> str:
    md = MarkItDown(enable_plugins=False)
    try:
        workbook = md.convert(to_bytesio(file))
    except (
        BadZipFile,
        ValueError,
        FileConversionException,
        UnsupportedFormatException,
    ) as e:
        error_str = f"Failed to extract text from {file_name or 'xlsx file'}: {e}"
        if file_name.startswith("~"):
            logger.debug(error_str + " (this is expected for files with ~)")
        else:
            logger.warning(error_str)
        return ""

    return workbook.markdown


def eml_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    text_file = io.TextIOWrapper(file, encoding=encoding)
    parser = EmailParser()
    message = parser.parse(text_file)

    text_content = []
    for part in message.walk():
        if part.get_content_type().startswith("text/plain"):
            payload = part.get_payload()
            if isinstance(payload, str):
                text_content.append(payload)
            elif isinstance(payload, list):
                text_content.extend(item for item in payload if isinstance(item, str))
            else:
                logger.warning(f"Unexpected payload type: {type(payload)}")
    return TEXT_SECTION_SEPARATOR.join(text_content)


def epub_to_text(file: IO[Any]) -> str:
    with zipfile.ZipFile(file) as epub:
        text_content = []
        for item in epub.infolist():
            if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                with epub.open(item) as html_file:
                    text_content.append(parse_html_page_basic(html_file))
        return TEXT_SECTION_SEPARATOR.join(text_content)


def file_io_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    file_content, _ = read_text_file(file, encoding=encoding)
    return file_content


def extract_file_text(
    file: IO[Any],
    file_name: str,
    break_on_unprocessable: bool = True,
    extension: str | None = None,
) -> str:
    """
    Legacy function that returns *only text*, ignoring embedded images.
    For backward-compatibility in code that only wants text.

    NOTE: Ignoring seems to be defined as returning an empty string for files it can't
    handle (such as images).
    """
    extension_to_function: dict[str, Callable[[IO[Any]], str]] = {
        ".pdf": pdf_to_text,
        ".docx": lambda f: docx_to_text_and_images(f, file_name)[0],  # no images
        ".pptx": lambda f: pptx_to_text(f, file_name),
        ".xlsx": lambda f: xlsx_to_text(f, file_name),
        ".eml": eml_to_text,
        ".epub": epub_to_text,
        ".html": parse_html_page_basic,
    }

    try:
        if get_unstructured_api_key():
            try:
                return unstructured_to_text(file, file_name)
            except Exception as unstructured_error:
                logger.error(
                    f"Failed to process with Unstructured: {str(unstructured_error)}. "
                    "Falling back to normal processing."
                )
        if extension is None:
            extension = get_file_ext(file_name)

        if is_accepted_file_ext(
            extension, OnyxExtensionType.Plain | OnyxExtensionType.Document
        ):
            func = extension_to_function.get(extension, file_io_to_text)
            file.seek(0)
            return func(file)

        # If unknown extension, maybe it's a text file
        file.seek(0)
        if is_text_file(file):
            return file_io_to_text(file)

        raise ValueError("Unknown file extension or not recognized as text data")

    except Exception as e:
        if break_on_unprocessable:
            raise RuntimeError(
                f"Failed to process file {file_name or 'Unknown'}: {str(e)}"
            ) from e
        logger.warning(f"Failed to process file {file_name or 'Unknown'}: {str(e)}")
        return ""


class ExtractionResult(NamedTuple):
    """Structured result from text and image extraction from various file types."""

    text_content: str
    embedded_images: Sequence[tuple[bytes, str]]
    metadata: dict[str, Any]


def extract_result_from_text_file(file: IO[Any]) -> ExtractionResult:
    encoding = detect_encoding(file)
    text_content_raw, file_metadata = read_text_file(
        file, encoding=encoding, ignore_onyx_metadata=False
    )
    return ExtractionResult(
        text_content=text_content_raw,
        embedded_images=[],
        metadata=file_metadata,
    )


def extract_text_and_images(
    file: IO[Any],
    file_name: str,
    pdf_pass: str | None = None,
    content_type: str | None = None,
) -> ExtractionResult:
    """
    Primary new function for the updated connector.
    Returns structured extraction result with text content, embedded images, and metadata.
    """
    file.seek(0)

    if get_unstructured_api_key():
        try:
            text_content = unstructured_to_text(file, file_name)
            return ExtractionResult(
                text_content=text_content, embedded_images=[], metadata={}
            )
        except Exception as e:
            logger.error(
                f"Failed to process with Unstructured: {str(e)}. "
                "Falling back to normal processing."
            )
            file.seek(0)  # Reset file pointer just in case

    # When we upload a document via a connector or MyDocuments, we extract and store the content of files
    # with content types in UploadMimeTypes.DOCUMENT_MIME_TYPES as plain text files.
    # As a result, the file name extension may differ from the original content type.
    # We process files with a plain text content type first to handle this scenario.
    if content_type == TEXT_MIME_TYPE:
        return extract_result_from_text_file(file)

    # Default processing
    try:
        extension = get_file_ext(file_name)

        # docx example for embedded images
        if extension == ".docx":
            text_content, images = docx_to_text_and_images(file, file_name)
            return ExtractionResult(
                text_content=text_content, embedded_images=images, metadata={}
            )

        # PDF example: we do not show complicated PDF image extraction here
        # so we simply extract text for now and skip images.
        if extension == ".pdf":
            text_content, pdf_metadata, images = read_pdf_file(
                file,
                pdf_pass,
                extract_images=get_image_extraction_and_analysis_enabled(),
            )
            return ExtractionResult(
                text_content=text_content, embedded_images=images, metadata=pdf_metadata
            )

        # For PPTX, XLSX, EML, etc., we do not show embedded image logic here.
        # You can do something similar to docx if needed.
        if extension == ".pptx":
            return ExtractionResult(
                text_content=pptx_to_text(file, file_name=file_name),
                embedded_images=[],
                metadata={},
            )

        if extension == ".xlsx":
            return ExtractionResult(
                text_content=xlsx_to_text(file, file_name=file_name),
                embedded_images=[],
                metadata={},
            )

        if extension == ".eml":
            return ExtractionResult(
                text_content=eml_to_text(file), embedded_images=[], metadata={}
            )

        if extension == ".epub":
            return ExtractionResult(
                text_content=epub_to_text(file), embedded_images=[], metadata={}
            )

        if extension == ".html":
            return ExtractionResult(
                text_content=parse_html_page_basic(file),
                embedded_images=[],
                metadata={},
            )

        # If we reach here and it's a recognized text extension
        if is_text_file_extension(file_name):
            return extract_result_from_text_file(file)

        # If it's an image file or something else, we do not parse embedded images from them
        # just return empty text
        return ExtractionResult(text_content="", embedded_images=[], metadata={})

    except Exception as e:
        logger.exception(f"Failed to extract text/images from {file_name}: {e}")
        return ExtractionResult(text_content="", embedded_images=[], metadata={})


def docx_to_txt_filename(file_path: str) -> str:
    return file_path.rsplit(".", 1)[0] + ".txt"
